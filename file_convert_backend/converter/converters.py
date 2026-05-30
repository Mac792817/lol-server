"""File conversion logic for images and documents."""
import os
import uuid
from io import BytesIO
from typing import Optional

from PIL import Image
from docx import Document
from openpyxl import load_workbook, Workbook
from pptx import Presentation


class ImageConverter:
    """Handles image format conversion."""

    SUPPORTED_FORMATS = ["jpg", "jpeg", "png", "webp", "gif", "bmp"]

    @classmethod
    def convert(
        cls,
        input_path: str,
        output_path: str,
        target_format: str,
        quality: int = 95
    ) -> bool:
        """Convert image to target format."""
        if target_format.lower() not in cls.SUPPORTED_FORMATS:
            raise ValueError(f"Unsupported target format: {target_format}")

        try:
            with Image.open(input_path) as img:
                if img.mode == "P" and img.format == "GIF":
                    img = img.convert("RGBA")

                if target_format.lower() in ["jpg", "jpeg"]:
                    if img.mode in ("RGBA", "LA", "P"):
                        background = Image.new("RGB", img.size, (255, 255, 255))
                        if img.mode == "P":
                            img = img.convert("RGBA")
                        if img.mode == "RGBA":
                            background.paste(img, mask=img.split()[-1])
                            img = background
                        else:
                            img = img.convert("RGB")
                    elif img.mode != "RGB":
                        img = img.convert("RGB")

                if target_format.lower() == "webp":
                    img.save(output_path, "WEBP", quality=quality)
                elif target_format.lower() in ["jpg", "jpeg"]:
                    img.save(output_path, "JPEG", quality=quality, optimize=True)
                elif target_format.lower() == "png":
                    img.save(output_path, "PNG", optimize=True)
                elif target_format.lower() == "gif":
                    img.save(output_path, "GIF")
                elif target_format.lower() == "bmp":
                    img.save(output_path, "BMP")
                else:
                    img.save(output_path)

                return True
        except Exception as e:
            raise RuntimeError(f"Image conversion failed: {str(e)}")

    @classmethod
    def get_image_info(cls, file_path: str) -> dict:
        """Get image metadata."""
        with Image.open(file_path) as img:
            return {
                "format": img.format,
                "mode": img.mode,
                "size": img.size,
                "width": img.width,
                "height": img.height,
            }


class DocumentConverter:
    """Handles document format conversion."""

    SUPPORTED_DOC_FORMATS = ["txt", "docx", "pdf"]
    SUPPORTED_SHEET_FORMATS = ["xlsx", "csv"]
    SUPPORTED_PRES_FORMATS = ["pptx"]

    @classmethod
    def txt_to_docx(cls, input_path: str, output_path: str) -> bool:
        """Convert TXT to DOCX."""
        doc = Document()
        with open(input_path, "r", encoding="utf-8") as f:
            content = f.read()
        doc.add_paragraph(content)
        doc.save(output_path)
        return True

    @classmethod
    def docx_to_txt(cls, input_path: str, output_path: str) -> bool:
        """Convert DOCX to TXT."""
        doc = Document(input_path)
        with open(output_path, "w", encoding="utf-8") as f:
            for para in doc.paragraphs:
                f.write(para.text + "\n")
        return True

    @classmethod
    def xlsx_to_csv(cls, input_path: str, output_path: str) -> bool:
        """Convert XLSX to CSV."""
        workbook = load_workbook(input_path, data_only=True)
        sheet = workbook.active

        with open(output_path, "w", encoding="utf-8", newline="") as f:
            for row in sheet.iter_rows(values_only=True):
                row_data = []
                for cell in row:
                    if cell is None:
                        row_data.append("")
                    else:
                        row_data.append(str(cell))
                f.write(",".join(row_data) + "\n")
        return True

    @classmethod
    def csv_to_xlsx(cls, input_path: str, output_path: str) -> bool:
        """Convert CSV to XLSX."""
        workbook = Workbook()
        sheet = workbook.active

        with open(input_path, "r", encoding="utf-8") as f:
            for row_idx, line in enumerate(f, start=1):
                values = line.strip().split(",")
                for col_idx, value in enumerate(values, start=1):
                    sheet.cell(row=row_idx, column=col_idx, value=value)

        workbook.save(output_path)
        return True

    @classmethod
    def pptx_to_txt(cls, input_path: str, output_path: str) -> bool:
        """Convert PPTX to TXT."""
        prs = Presentation(input_path)
        with open(output_path, "w", encoding="utf-8") as f:
            for slide_num, slide in enumerate(prs.slides, start=1):
                f.write(f"=== Slide {slide_num} ===\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        f.write(shape.text + "\n")
                f.write("\n")
        return True


def get_converter(file_format: str):
    """Get appropriate converter based on file format."""
    file_format = file_format.lower()

    if file_format in ImageConverter.SUPPORTED_FORMATS:
        return ImageConverter
    elif file_format in DocumentConverter.SUPPORTED_DOC_FORMATS:
        return DocumentConverter
    elif file_format in DocumentConverter.SUPPORTED_SHEET_FORMATS:
        return DocumentConverter
    elif file_format in DocumentConverter.SUPPORTED_PRES_FORMATS:
        return DocumentConverter

    return None


def is_image_format(file_format: str) -> bool:
    """Check if format is an image format."""
    return file_format.lower() in ImageConverter.SUPPORTED_FORMATS


def is_document_format(file_format: str) -> bool:
    """Check if format is a document format."""
    file_format = file_format.lower()
    return (
        file_format in DocumentConverter.SUPPORTED_DOC_FORMATS
        or file_format in DocumentConverter.SUPPORTED_SHEET_FORMATS
        or file_format in DocumentConverter.SUPPORTED_PRES_FORMATS
    )
